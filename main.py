import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from lyricsgenius import Genius
import subprocess

GENIUS_TOKEN = ""

prs = Presentation("template.pptx")
lyricsLayout = prs.slide_layouts[1]

if GENIUS_TOKEN == "":
    print("No genius API key found. Set the GENIUS_TOKEN variable with a valid Genius API key.")
else:
    genius = Genius(GENIUS_TOKEN)

def addSlide(lyrics, fontSize = 40, fontWeight = "Medium"):
    slide = prs.slides.add_slide(lyricsLayout)
    shapes = slide.shapes
    for shape in shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = shape.text_frame.paragraphs[0]

            paragraph.text = lyrics

            font = paragraph.font
            font.size = Pt(fontSize)
            font.name = "Poppins " + fontWeight
            font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.line_spacing = 1

def songSearch():
    searchTerm = input("What song would you like to generate slides for? ")
    print("Searching...")
    songs = genius.search_songs(searchTerm)
    desiredSong = {}

    for song in songs['hits']:
        print("Is this the song you were looking for?", song['result']['full_title'])
        if input("y or n: ") in ("y", "Y", "yes", "Yes"):
            print("Great! Getting lyrics...")
            return genius.search_song(song_id=song['result']['id'])
    
    print("Sorry, your song wasn't found")
    return None

def tidyLyrics(lyrics):
    lines = []
    headers = ("Chorus]", "[Verse", "[Outro]", "[Bridge", "[Tag]", "Bridge]", "[Refrain]")

    for line in lyrics.splitlines():
        if any(header in line for header in headers):
            if lines and lines[-1]:
                lines.append("")
        else:
            lines.append(line)

    output = ""
    for line in lines:
        output += line + "\n"

    return output

def editLyrics(lyrics):
    songFile = "editLyrics.txt"
    text_file = open(songFile, "w")
    n = text_file.write(tidyLyrics(lyrics))
    text_file.close()

    subprocess.call(['open', '-a', 'TextEdit', songFile])
    input("Edit the file as desired. Line breaks mark where slides should be separated. Press enter to continue")
    print("\n")
    return songFile

def readLyrics(filename):
    lyrics = [""]
    with open(filename) as f:
        for line in f.readlines():
            if line.isspace():
                lyrics.append("")
            else:
                lyrics[-1] += line
    return lyrics

def makeSongSlides(lyrics):
    for slideLines in lyrics:
        addSlide(slideLines.strip())

def main():
    if GENIUS_TOKEN == "": # if no valid API key quit
        return

    songFile = ""
    songName = ""

    while True:
        print("Do you have a lyric file, or would you like to search for a song?")
        print("1) use lyric file")
        print("2) search for a song")
        useFile = input("Enter 1 or 2: ")
        if useFile == "1":
            songFile = input("Enter the path to your lyrics file: ")
            songName = input("Enter the name of this song: ")
            break
        if useFile == "2":
            song = songSearch()
            songFile = editLyrics(tidyLyrics(song.lyrics))
            songName = song.title
            break
        print("Invalid choice\n")
        

    lyrics = readLyrics(songFile)
    addSlide(songName, fontSize=60, fontWeight="SemiBold") # add title slide
    makeSongSlides(lyrics)
    outFile = "output/" + songName + ".pptx"
    prs.save(outFile)
    print("Your slides have been created! The file has been saved to " + outFile)

if __name__ == "__main__":
    main()